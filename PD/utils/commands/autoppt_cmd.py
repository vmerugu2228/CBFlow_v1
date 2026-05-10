#!/usr/bin/env python3
"""
CBflow AutoPPT — Automatic PowerPoint / HTML Summary Generator
Matches PD Run Summary Template layout:
  Slide 1: Summary — 19-column QoR table + image placeholders + metadata
  Slide 2: Place — Timing QoR + images + summary
  Slide 3: CTS — Timing QoR + clock skew table + summary
  Slide 4: ctsOpt — Timing QOR + images + summary
  Slide 5: routeOpt — Timing QOR + images + summary

Uses python-pptx if available; falls back to HTML slide deck.
"""

import argparse
import json
import os
import re
import sys
from datetime import datetime
from pathlib import Path
from collections import OrderedDict

try:
    from .logging_config import get_logger
    logger = get_logger('cbflow.autoppt')
except Exception:
    import logging
    logger = logging.getLogger('cbflow.autoppt')

# ── python-pptx availability ────────────────────────────────────────────────
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False

# ── Constants (matching template theme) ──────────────────────────────────────
ACCENT1 = '#4472C4'   # Blue headers
ACCENT2 = '#ED7D31'   # Orange highlights
ACCENT3 = '#A5A5A5'   # Gray borders
GREEN   = '#70AD47'
RED     = '#FF0000'
DK2     = '#44546A'    # Dark blue-gray
LT2     = '#E7E6E6'    # Light gray bg
WHITE   = '#FFFFFF'
BLACK   = '#000000'

# Stage mapping: CBflow node name → display name in PPT
STAGE_MAP = OrderedDict([
    ('init_design1', 'Initial'),
    ('synthesis1',   'compile'),
    ('place1',       'place'),
    ('cts1',         'cts'),
    ('cts_opt1',     'ctsOpt'),
    ('route1',       'route'),
    ('pro1',         'routeOpt'),
    ('signoff1',     'chipFinish'),
])

# FC block names used in reports
BLOCK_MAP = {
    'init_design1': 'init_design',
    'synthesis1':   'compile',
    'place1':       'place_opt',
    'cts1':         'clock_opt_cts',
    'cts_opt1':     'clock_opt_opto',
    'route1':       'route_auto',
    'pro1':         'route_opt',
    'signoff1':     'chip_finish',
}


# ══════════════════════════════════════════════════════════════════════════════
# DATA COLLECTION
# ══════════════════════════════════════════════════════════════════════════════

def collect_ppt_data(run_dir: str) -> dict:
    """Collect all data needed for PPT generation."""
    data = {
        'run_name': os.path.basename(run_dir),
        'run_dir': run_dir,
        'user': os.environ.get('USER', 'unknown'),
        'timestamp': datetime.now().strftime('%Y-%m-%d %H:%M:%S'),
        'flow_type': '', 'design_name': '', 'project_name': '',
        'tech_node': '', 'release_version': '',
        'tile_name': '', 'tile_owner': '',
        'netlist_release': '', 'fp_release': '',
        'tool_version': '', 'setup_corner': '', 'hold_corner': '',
        # Per-stage QoR: stage -> {wns_setup, tns_setup, fep_setup, wns_hold, ...}
        'stage_qor': OrderedDict(),
        # Clock table: list of {clock_name, clock_id, global_skew}
        'clock_table': [],
        # Per-stage summary text
        'stage_summary': {},
        # Image paths per stage
        'stage_images': {},
        # Run times per stage
        'stage_runtime': {},
    }

    _parse_run_env(run_dir, data)
    _collect_stage_qor(run_dir, data)
    _collect_clock_data(run_dir, data)
    _collect_stage_runtimes(run_dir, data)
    _find_images(run_dir, data)

    return data


def _parse_run_env(run_dir: str, data: dict):
    """Parse .run.cbflow.tcl for run metadata."""
    env_file = os.path.join(run_dir, '.run.cbflow.tcl')
    if not os.path.exists(env_file):
        return
    try:
        with open(env_file) as f:
            for line in f:
                line = line.strip()
                for pattern, key in [
                    (r'set\s+flow\(design_name\)\s+"([^"]+)"', 'design_name'),
                    (r'set\s+project\(name\)\s+"([^"]+)"', 'project_name'),
                    (r'set\s+::env\(TECH_NAME\)\s+"([^"]+)"', 'tech_node'),
                    (r'set\s+::env\(CBFLOW_RELEASE\)\s+"([^"]+)"', 'release_version'),
                    (r'set\s+flow\(type\)\s+"([^"]+)"', 'flow_type'),
                ]:
                    m = re.match(pattern, line)
                    if m and not data.get(key):
                        data[key] = m.group(1)
    except Exception:
        pass

    if not data['flow_type']:
        m = re.match(r'P\d+_run_(\w+)_run\d+', data['run_name'])
        if m:
            data['flow_type'] = m.group(1)

    # Try to get tool version
    try:
        fc_version_file = Path(run_dir) / 'work' / data['flow_type'] / 'init_design1' / 'run' / 'fc_version.txt'
        if fc_version_file.exists():
            data['tool_version'] = fc_version_file.read_text().strip()
    except Exception:
        pass


def _collect_stage_qor(run_dir: str, data: dict):
    """Parse report_qor.rpt from each stage's reports directory."""
    flow_type = data.get('flow_type', 'SYNTH_PNR')
    work_dir = Path(run_dir) / 'work' / flow_type

    for node_name, display_name in STAGE_MAP.items():
        reports_dir = work_dir / node_name / 'reports'
        qor = {
            'wns_setup': '', 'tns_setup': '', 'fep_setup': '',
            'wns_hold': '', 'tns_hold': '', 'fep_hold': '',
            'wns_io_setup': '', 'tns_io_setup': '', 'fep_io_setup': '',
            'std_count': '', 'std_area': '', 'holdbuf_cnt': '',
            'clock_id': '', 'clock_buf_inv_cgc': '',
            'congestion': '', 'drc_shorts': '',
            'leakage_power': '', 'runtime': '',
        }

        # Parse report_qor.rpt
        qor_file = reports_dir / 'report_qor.rpt'
        if qor_file.exists():
            _parse_qor_file(qor_file, qor)

        # Parse report_qor_summary.rpt as fallback
        qor_summary = reports_dir / 'report_qor_summary.rpt'
        if qor_summary.exists() and not qor['wns_setup']:
            _parse_qor_file(qor_summary, qor)

        # Parse timing reports for WNS
        for suffix, prefix in [('max', 'setup'), ('min', 'hold')]:
            timing_file = reports_dir / f'report_timing.{suffix}.rpt'
            if timing_file.exists() and not qor[f'wns_{prefix}']:
                _parse_timing_file(timing_file, qor, prefix)

        # Parse power
        power_file = reports_dir / 'report_power.rpt'
        if power_file.exists() and not qor['leakage_power']:
            _parse_power_file(power_file, qor)

        # Parse congestion
        congestion_file = reports_dir / 'report_congestion.rpt'
        if congestion_file.exists() and not qor['congestion']:
            _parse_congestion_file(congestion_file, qor)

        # Parse DRC
        _parse_drc_from_reports(reports_dir, qor)

        # Parse design stats
        design_file = reports_dir / 'report_design.rpt'
        if design_file.exists():
            _parse_design_file(design_file, qor)

        utilization_file = reports_dir / 'report_utilization.rpt'
        if utilization_file.exists():
            _parse_utilization_file(utilization_file, qor)

        data['stage_qor'][display_name] = qor


def _parse_qor_file(filepath: Path, qor: dict):
    """Extract metrics from report_qor.rpt."""
    try:
        content = filepath.read_text(errors='ignore')
        patterns = [
            (r'WNS\s*\(?\s*setup\s*\)?\s*[:=]\s*([-\d.]+)', 'wns_setup'),
            (r'TNS\s*\(?\s*setup\s*\)?\s*[:=]\s*([-\d.]+)', 'tns_setup'),
            (r'(?:Failing\s+Endpoints|FEP)\s*\(?\s*setup\s*\)?\s*[:=]\s*(\d+)', 'fep_setup'),
            (r'WNS\s*\(?\s*hold\s*\)?\s*[:=]\s*([-\d.]+)', 'wns_hold'),
            (r'TNS\s*\(?\s*hold\s*\)?\s*[:=]\s*([-\d.]+)', 'tns_hold'),
            (r'(?:Failing\s+Endpoints|FEP)\s*\(?\s*hold\s*\)?\s*[:=]\s*(\d+)', 'fep_hold'),
            (r'Leaf\s+Cell\s+Count\s*[:=]\s*([\d,]+)', 'std_count'),
            (r'Cell\s+Area\s*(?:\([^)]*\))?\s*[:=]\s*([\d,.]+)', 'std_area'),
        ]
        for pattern, key in patterns:
            m = re.search(pattern, content, re.IGNORECASE)
            if m and not qor[key]:
                qor[key] = m.group(1).replace(',', '')
    except Exception:
        pass


def _parse_timing_file(filepath: Path, qor: dict, check_type: str):
    """Extract worst slack from timing report."""
    try:
        content = filepath.read_text(errors='ignore')
        # Look for slack line
        m = re.search(r'slack\s*\(?(MET|VIOLATED)\)?\s*([-\d.]+)', content)
        if m:
            qor[f'wns_{check_type}'] = m.group(2)
    except Exception:
        pass


def _parse_power_file(filepath: Path, qor: dict):
    """Extract leakage power."""
    try:
        content = filepath.read_text(errors='ignore')
        m = re.search(r'(?:Cell\s+)?Leakage\s+Power\s*[:=]\s*([\d.eE+-]+)\s*(\w+)', content, re.I)
        if m:
            qor['leakage_power'] = f"{m.group(1)} {m.group(2)}"
    except Exception:
        pass


def _parse_congestion_file(filepath: Path, qor: dict):
    """Extract congestion metrics."""
    try:
        content = filepath.read_text(errors='ignore')
        m = re.search(r'(?:Both|Overall).*?overflow.*?(\d+)', content, re.I)
        if m:
            qor['congestion'] = m.group(1)
    except Exception:
        pass


def _parse_drc_from_reports(reports_dir: Path, qor: dict):
    """Collect DRC error count."""
    for drc_file in reports_dir.glob('*drc*'):
        try:
            content = drc_file.read_text(errors='ignore')
            m = re.search(r'Total.*?(\d+)\s+(?:error|violation)', content, re.I)
            if m:
                qor['drc_shorts'] = m.group(1)
                return
        except Exception:
            pass


def _parse_design_file(filepath: Path, qor: dict):
    """Extract cell count from report_design."""
    try:
        content = filepath.read_text(errors='ignore')
        m = re.search(r'Number\s+of\s+(?:leaf\s+)?cells?\s*[:=]\s*([\d,]+)', content, re.I)
        if m and not qor['std_count']:
            qor['std_count'] = m.group(1).replace(',', '')
    except Exception:
        pass


def _parse_utilization_file(filepath: Path, qor: dict):
    """Extract area from report_utilization."""
    try:
        content = filepath.read_text(errors='ignore')
        m = re.search(r'Cell\s+Area\s*[:=]\s*([\d,.]+)', content, re.I)
        if m and not qor['std_area']:
            qor['std_area'] = m.group(1).replace(',', '')
    except Exception:
        pass


def _collect_clock_data(run_dir: str, data: dict):
    """Parse clock_qor report for CTS slide table."""
    flow_type = data.get('flow_type', 'SYNTH_PNR')
    # Try cts1 first, then cts_opt1
    for node in ['cts1', 'cts_opt1']:
        rpt = Path(run_dir) / 'work' / flow_type / node / 'reports' / 'report_clock_qor.rpt'
        if rpt.exists():
            try:
                content = rpt.read_text(errors='ignore')
                # Look for clock table rows: clock_name  skew  latency etc
                for m in re.finditer(r'^\s*(\S+)\s+(\S+)\s+([\d.]+)\s', content, re.M):
                    data['clock_table'].append({
                        'clock_name': m.group(1),
                        'clock_id': m.group(2),
                        'global_skew': m.group(3),
                    })
            except Exception:
                pass
            if data['clock_table']:
                break


def _collect_stage_runtimes(run_dir: str, data: dict):
    """Collect runtime from engine DB or stamps."""
    try:
        from status_provider import get_status_provider
        provider = get_status_provider(run_dir)
        all_status = provider.get_all_status()
        for stage, info in all_status.items():
            ts = info.get('timestamp', '')
            time_short = ts.split('T')[-1][:8] if 'T' in ts else ts[-8:] if len(ts) > 8 else ts
            data['stage_runtime'][stage] = time_short
    except Exception:
        stamps_dir = Path(run_dir) / '.stamps'
        if not stamps_dir.is_dir():
            return
        for stamp in sorted(stamps_dir.glob('*.stamp')):
            try:
                mtime = datetime.fromtimestamp(stamp.stat().st_mtime)
                data['stage_runtime'][stamp.stem] = mtime.strftime('%H:%M:%S')
            except Exception:
                pass


def _find_images(run_dir: str, data: dict):
    """Find GUI snapshot images for each stage."""
    flow_type = data.get('flow_type', 'SYNTH_PNR')
    work_dir = Path(run_dir) / 'work' / flow_type

    image_patterns = {
        'floorplan':    ['*floorplan*', '*fp*'],
        'cell_density': ['*cell_density*', '*density*'],
        'congestion':   ['*congestion*', '*cong*'],
        'pin_density':  ['*pin_density*', '*pin*'],
        'drc':          ['*drc*'],
        'placement':    ['*placement*', '*place*'],
    }

    for node_name in STAGE_MAP:
        images = {}
        img_dirs = [
            work_dir / node_name / 'reports',
            work_dir / node_name / 'images',
            work_dir / node_name / 'snapshots',
        ]
        for img_dir in img_dirs:
            if not img_dir.is_dir():
                continue
            for img_type, patterns in image_patterns.items():
                if img_type in images:
                    continue
                for pat in patterns:
                    for ext in ['*.png', '*.jpg', '*.jpeg', '*.gif']:
                        for found in img_dir.glob(f'{pat}{ext}'):
                            images[img_type] = str(found)
                            break
                    if img_type in images:
                        break
        if images:
            data['stage_images'][node_name] = images


# ══════════════════════════════════════════════════════════════════════════════
# PPTX GENERATION (python-pptx)
# ══════════════════════════════════════════════════════════════════════════════

def generate_pptx(data: dict, output_path: str):
    """Generate PPTX matching the PD Run Summary Template."""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    _pptx_slide1_summary(prs, data)
    _pptx_stage_slide(prs, data, 'place1', 'Place',
                      ['placement', 'cell_density', 'congestion', 'pin_density'])
    _pptx_slide_cts(prs, data)
    _pptx_stage_slide(prs, data, 'cts_opt1', 'ctsOpt',
                      ['placement', 'cell_density', 'congestion', 'pin_density'])
    _pptx_stage_slide(prs, data, 'pro1', 'routeOpt',
                      ['drc', 'cell_density', 'congestion', 'pin_density'])

    prs.save(output_path)


def _rgb(hex_color: str) -> RGBColor:
    """Convert hex color to RGBColor."""
    h = hex_color.lstrip('#')
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def _set_cell(table, row, col, text, font_size=8, bold=False, color=BLACK, align=PP_ALIGN.CENTER):
    """Set table cell text with formatting."""
    cell = table.cell(row, col)
    cell.text = str(text)
    for p in cell.text_frame.paragraphs:
        p.font.size = Pt(font_size)
        p.font.bold = bold
        p.font.color.rgb = _rgb(color)
        p.alignment = align


def _add_textbox(slide, left, top, width, height, text, size=12, bold=False, color=BLACK):
    """Add a text box to a slide."""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = _rgb(color)
    return txBox


def _add_image_placeholder(slide, left, top, width, height, label, image_path=None):
    """Add image or placeholder rectangle."""
    if image_path and os.path.exists(image_path):
        slide.shapes.add_picture(image_path, Inches(left), Inches(top), Inches(width), Inches(height))
    else:
        shape = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(width), Inches(height))
        shape.fill.solid()
        shape.fill.fore_color.rgb = _rgb(LT2)
        shape.line.color.rgb = _rgb(ACCENT3)
        tf = shape.text_frame
        p = tf.paragraphs[0]
        p.text = label
        p.font.size = Pt(9)
        p.font.color.rgb = _rgb(ACCENT3)
        p.alignment = PP_ALIGN.CENTER


def _pptx_slide1_summary(prs, data):
    """Slide 1: Summary — main QoR table matching template."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    # Title
    _add_textbox(slide, 0.2, 0.1, 7, 0.4, f"Run : {data['run_name']} Summary :", 14, True, DK2)

    # Metadata table (top-right)
    meta_rows = [
        ('Tile Name', data.get('design_name', '')),
        ('Tile Owner', data.get('user', '')),
        ('Netlist Release', data.get('netlist_release', '')),
        ('Floorplan Release', data.get('fp_release', '')),
        ('Tool Version', data.get('tool_version', '')),
        ('CBflow Release', data.get('release_version', '')),
    ]
    meta_tbl = slide.shapes.add_table(len(meta_rows), 2, Inches(9.4), Inches(0.1), Inches(3.8), Inches(1.6)).table
    for r, (k, v) in enumerate(meta_rows):
        _set_cell(meta_tbl, r, 0, k, 7, True, DK2, PP_ALIGN.LEFT)
        _set_cell(meta_tbl, r, 1, v, 7, False, BLACK, PP_ALIGN.LEFT)

    # ── Main 19-column QoR table ──
    stages = list(STAGE_MAP.values())
    n_data_rows = len(stages)
    total_rows = 4 + n_data_rows  # 2 header merge rows + col headers + sub-headers + data
    n_cols = 19

    tbl_shape = slide.shapes.add_table(total_rows, n_cols, Inches(0.1), Inches(1.0), Inches(13.1), Inches(0.22 * total_rows))
    table = tbl_shape.table

    # Row 0: Run Name header (merged)
    _set_cell(table, 0, 0, '', 7)
    table.cell(0, 1).merge(table.cell(0, 18))
    _set_cell(table, 0, 1, f'Run Name : {data["run_name"]}', 9, True, WHITE)
    table.cell(0, 1).fill.solid()
    table.cell(0, 1).fill.fore_color.rgb = _rgb(ACCENT1)

    # Row 1: Setup Corner / Hold Corner (merged)
    _set_cell(table, 1, 0, '', 7)
    table.cell(1, 1).merge(table.cell(1, 10))
    _set_cell(table, 1, 1, f'Setup Corner : {data.get("setup_corner", "")}', 8, True, DK2)
    table.cell(1, 11).merge(table.cell(1, 18))
    _set_cell(table, 1, 11, f'Hold Corner : {data.get("hold_corner", "")}', 8, True, DK2)

    # Row 2: Category headers (merged)
    cat_headers = [
        ('', 1), ('Setup R2R', 3), ('Hold', 3), ('IO Setup', 3),
        ('Design Stats', 3), ('Clock Stats', 2), ('Routability', 2),
        ('Power', 1), ('Run Stats', 1),
    ]
    col = 0
    for label, span in cat_headers:
        if span > 1:
            table.cell(2, col).merge(table.cell(2, col + span - 1))
        _set_cell(table, 2, col, label, 7, True, WHITE)
        table.cell(2, col).fill.solid()
        table.cell(2, col).fill.fore_color.rgb = _rgb(DK2)
        col += span

    # Row 3: Column headers
    col_headers = ['Node', 'WNS', 'TNS', 'FEP', 'WNS', 'TNS', 'FEP',
                   'WNS', 'TNS', 'FEP', 'Std. Count', 'Std. Area', 'Holdbuf.cnt',
                   'Clock ID', 'Clock Buf/Inv/Cgc', 'Congestion', 'DRC/Shorts', 'Lkg', 'Run Time']
    for c, hdr in enumerate(col_headers):
        _set_cell(table, 3, c, hdr, 6, True, DK2)
        table.cell(3, c).fill.solid()
        table.cell(3, c).fill.fore_color.rgb = _rgb(LT2)

    # Data rows
    for r, stage_name in enumerate(stages, 4):
        qor = data['stage_qor'].get(stage_name, {})
        vals = [
            stage_name,
            qor.get('wns_setup', ''), qor.get('tns_setup', ''), qor.get('fep_setup', ''),
            qor.get('wns_hold', ''), qor.get('tns_hold', ''), qor.get('fep_hold', ''),
            qor.get('wns_io_setup', ''), qor.get('tns_io_setup', ''), qor.get('fep_io_setup', ''),
            qor.get('std_count', ''), qor.get('std_area', ''), qor.get('holdbuf_cnt', ''),
            qor.get('clock_id', ''), qor.get('clock_buf_inv_cgc', ''),
            qor.get('congestion', ''), qor.get('drc_shorts', ''),
            qor.get('leakage_power', ''), qor.get('runtime', ''),
        ]
        for c, v in enumerate(vals):
            color = BLACK
            # Color negative WNS red
            if c in (1, 4, 7) and v and v.startswith('-'):
                color = RED
            _set_cell(table, r, c, v, 7, c == 0, color)

    # Summary text box
    _add_textbox(slide, 0.3, 3.3, 7.2, 1.8, 'Summary : Issues : Jira Tickets:', 10, False, DK2)

    # Image placeholders (2x2 grid, right side)
    images = {}
    for node in STAGE_MAP:
        if node in data.get('stage_images', {}):
            images.update(data['stage_images'][node])
    _add_image_placeholder(slide, 7.7, 3.4, 2.6, 1.9, 'Floorplan Image', images.get('floorplan'))
    _add_image_placeholder(slide, 10.5, 3.4, 2.6, 1.9, 'Place Cell Density', images.get('cell_density'))
    _add_image_placeholder(slide, 7.7, 5.5, 2.6, 1.9, 'Place Congestion', images.get('congestion'))
    _add_image_placeholder(slide, 10.5, 5.5, 2.6, 1.9, 'Route DRC', images.get('drc'))


def _pptx_stage_slide(prs, data, node_name, display_name, image_types):
    """Generic per-stage slide: Title + Timing QoR + Summary + 4 images."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Stage title
    _add_textbox(slide, 0.2, 0.1, 4, 0.4, f'{display_name} :', 14, True, DK2)

    # Timing QoR section (left)
    _add_textbox(slide, 0.4, 0.7, 3, 0.3, 'Timing QOR', 11, True, ACCENT1)

    qor = data['stage_qor'].get(display_name, {})
    timing_text = f"""Setup WNS: {qor.get('wns_setup', 'N/A')}   TNS: {qor.get('tns_setup', 'N/A')}   FEP: {qor.get('fep_setup', 'N/A')}
Hold  WNS: {qor.get('wns_hold', 'N/A')}   TNS: {qor.get('tns_hold', 'N/A')}   FEP: {qor.get('fep_hold', 'N/A')}
Std Cells: {qor.get('std_count', 'N/A')}   Area: {qor.get('std_area', 'N/A')}
Congestion: {qor.get('congestion', 'N/A')}   DRC: {qor.get('drc_shorts', 'N/A')}
Leakage: {qor.get('leakage_power', 'N/A')}"""

    _add_textbox(slide, 0.4, 1.1, 7, 2.0, timing_text, 9, False, BLACK)

    # Summary section (right)
    _add_textbox(slide, 7.7, 0.7, 5, 0.3, 'Summary:', 11, True, ACCENT1)
    _add_textbox(slide, 7.7, 1.1, 5.4, 2.0,
                 data.get('stage_summary', {}).get(display_name, ''), 9, False, DK2)

    # Image placeholders (2x2 grid)
    img_labels = {
        'placement':    f'{display_name} Placement',
        'cell_density': f'{display_name} Cell Density',
        'congestion':   f'{display_name} Congestion',
        'pin_density':  f'{display_name} Pin Density',
        'drc':          f'{display_name} DRC',
    }
    images = data.get('stage_images', {}).get(node_name, {})
    positions = [(7.7, 3.4), (10.5, 3.4), (7.7, 5.5), (10.5, 5.5)]
    for i, img_type in enumerate(image_types[:4]):
        x, y = positions[i]
        label = img_labels.get(img_type, img_type)
        _add_image_placeholder(slide, x, y, 2.6, 1.9, label, images.get(img_type))


def _pptx_slide_cts(prs, data):
    """CTS slide: Timing QoR + Clock skew table."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    _add_textbox(slide, 0.2, 0.1, 4, 0.4, 'CTS :', 14, True, DK2)
    _add_textbox(slide, 0.4, 0.7, 3, 0.3, 'Timing QOR', 11, True, ACCENT1)

    qor = data['stage_qor'].get('cts', {})
    timing_text = f"""Setup WNS: {qor.get('wns_setup', 'N/A')}   TNS: {qor.get('tns_setup', 'N/A')}   FEP: {qor.get('fep_setup', 'N/A')}
Hold  WNS: {qor.get('wns_hold', 'N/A')}   TNS: {qor.get('tns_hold', 'N/A')}   FEP: {qor.get('fep_hold', 'N/A')}"""
    _add_textbox(slide, 0.4, 1.1, 7, 1.0, timing_text, 9, False, BLACK)

    _add_textbox(slide, 7.7, 0.7, 5, 0.3, 'Summary:', 11, True, ACCENT1)

    # Clock skew table
    clock_rows = data.get('clock_table', [])
    n_rows = max(len(clock_rows), 1) + 1  # header + data
    n_rows = min(n_rows, 22)  # cap at 22 like template

    tbl_shape = slide.shapes.add_table(n_rows, 3, Inches(0.4), Inches(2.2), Inches(7.0), Inches(0.25 * n_rows))
    table = tbl_shape.table

    for c, hdr in enumerate(['Clock Name', 'Clock ID', 'Global Skew']):
        _set_cell(table, 0, c, hdr, 8, True, WHITE)
        table.cell(0, c).fill.solid()
        table.cell(0, c).fill.fore_color.rgb = _rgb(ACCENT1)

    for r, clk in enumerate(clock_rows[:n_rows - 1], 1):
        _set_cell(table, r, 0, clk.get('clock_name', ''), 7, False, BLACK, PP_ALIGN.LEFT)
        _set_cell(table, r, 1, clk.get('clock_id', ''), 7, False, BLACK)
        _set_cell(table, r, 2, clk.get('global_skew', ''), 7, False, BLACK)


# ══════════════════════════════════════════════════════════════════════════════
# HTML FALLBACK
# ══════════════════════════════════════════════════════════════════════════════

def generate_html_slides(data: dict, output_path: str):
    """Generate HTML slide deck matching the PD Run Summary Template layout."""
    stages = list(STAGE_MAP.values())

    # ── Slide 1: Summary QoR table ──
    qor_header = """<tr class="cat-row">
      <td></td><td colspan="3">Setup R2R</td><td colspan="3">Hold</td>
      <td colspan="3">IO Setup</td><td colspan="3">Design Stats</td>
      <td colspan="2">Clock Stats</td><td colspan="2">Routability</td>
      <td>Power</td><td>Run Stats</td></tr>
    <tr class="hdr-row">
      <th>Node</th><th>WNS</th><th>TNS</th><th>FEP</th>
      <th>WNS</th><th>TNS</th><th>FEP</th>
      <th>WNS</th><th>TNS</th><th>FEP</th>
      <th>Std.Cnt</th><th>Std.Area</th><th>HoldBuf</th>
      <th>ClkID</th><th>Buf/Inv/Cgc</th>
      <th>Cong</th><th>DRC</th><th>Lkg</th><th>Runtime</th></tr>"""

    qor_rows = ""
    for s in stages:
        q = data['stage_qor'].get(s, {})
        def _td(val, is_wns=False):
            cls = ' class="neg"' if is_wns and str(val).startswith('-') else ''
            return f'<td{cls}>{val}</td>'
        qor_rows += f"""<tr><td class="stage">{s}</td>
          {_td(q.get('wns_setup',''), True)}{_td(q.get('tns_setup',''))}{_td(q.get('fep_setup',''))}
          {_td(q.get('wns_hold',''), True)}{_td(q.get('tns_hold',''))}{_td(q.get('fep_hold',''))}
          {_td(q.get('wns_io_setup',''))}{_td(q.get('tns_io_setup',''))}{_td(q.get('fep_io_setup',''))}
          {_td(q.get('std_count',''))}{_td(q.get('std_area',''))}{_td(q.get('holdbuf_cnt',''))}
          {_td(q.get('clock_id',''))}{_td(q.get('clock_buf_inv_cgc',''))}
          {_td(q.get('congestion',''))}{_td(q.get('drc_shorts',''))}
          {_td(q.get('leakage_power',''))}{_td(q.get('runtime',''))}</tr>\n"""

    # Metadata
    meta = f"""<table class="meta"><tr><td class="lbl">Tile Name</td><td>{data.get('design_name','')}</td></tr>
      <tr><td class="lbl">Owner</td><td>{data.get('user','')}</td></tr>
      <tr><td class="lbl">Tool Version</td><td>{data.get('tool_version','')}</td></tr>
      <tr><td class="lbl">CBflow Release</td><td>{data.get('release_version','')}</td></tr>
      <tr><td class="lbl">Technology</td><td>{data.get('tech_node','')}</td></tr></table>"""

    # ── Per-stage slides ──
    stage_slides = ""
    for node_name, display_name in [('place1','Place'), ('cts1','CTS'), ('cts_opt1','ctsOpt'), ('pro1','routeOpt')]:
        q = data['stage_qor'].get(display_name, {})
        clock_tbl = ""
        if display_name == 'CTS' and data.get('clock_table'):
            clock_tbl = '<table class="clock-tbl"><tr><th>Clock Name</th><th>Clock ID</th><th>Skew</th></tr>'
            for clk in data['clock_table'][:20]:
                clock_tbl += f"<tr><td>{clk.get('clock_name','')}</td><td>{clk.get('clock_id','')}</td><td>{clk.get('global_skew','')}</td></tr>"
            clock_tbl += '</table>'

        stage_slides += f"""
    <div class="slide">
      <div class="slide-hdr"><span class="stage-title">{display_name}</span></div>
      <div class="two-col">
        <div class="left-col">
          <h3>Timing QOR</h3>
          <table class="qor-mini">
            <tr><td>Setup WNS</td><td class="{'neg' if str(q.get('wns_setup','')).startswith('-') else ''}">{q.get('wns_setup','N/A')}</td>
                <td>TNS</td><td>{q.get('tns_setup','N/A')}</td><td>FEP</td><td>{q.get('fep_setup','N/A')}</td></tr>
            <tr><td>Hold WNS</td><td>{q.get('wns_hold','N/A')}</td>
                <td>TNS</td><td>{q.get('tns_hold','N/A')}</td><td>FEP</td><td>{q.get('fep_hold','N/A')}</td></tr>
            <tr><td>Std Cells</td><td>{q.get('std_count','N/A')}</td><td>Area</td><td>{q.get('std_area','N/A')}</td>
                <td>Cong</td><td>{q.get('congestion','N/A')}</td></tr>
            <tr><td>DRC</td><td>{q.get('drc_shorts','N/A')}</td><td>Lkg</td><td colspan="3">{q.get('leakage_power','N/A')}</td></tr>
          </table>
          {clock_tbl}
        </div>
        <div class="right-col">
          <h3>Summary</h3>
          <div class="summary-text">{data.get('stage_summary',{}).get(display_name,'')}</div>
          <div class="img-grid">
            <div class="img-ph">Cell Density</div><div class="img-ph">Congestion</div>
            <div class="img-ph">Pin Density</div><div class="img-ph">DRC</div>
          </div>
        </div>
      </div>
    </div>"""

    html = f"""<!DOCTYPE html>
<html><head><meta charset="utf-8"><title>PD Run Summary - {data['run_name']}</title>
<style>
  * {{ box-sizing: border-box; margin: 0; padding: 0; }}
  body {{ font-family: Calibri, 'Segoe UI', Arial, sans-serif; background: #e8e8e8; }}
  .slide {{ background: white; width: 1000px; margin: 20px auto; padding: 20px 25px;
           border-radius: 4px; box-shadow: 0 2px 8px rgba(0,0,0,0.15);
           min-height: 560px; page-break-after: always; position: relative; }}
  .slide-hdr {{ border-bottom: 3px solid {ACCENT1}; margin-bottom: 12px; padding-bottom: 4px; }}
  .stage-title {{ font-size: 16px; font-weight: bold; color: {DK2}; }}
  h3 {{ color: {ACCENT1}; font-size: 13px; margin: 8px 0 6px; }}
  .meta {{ position: absolute; right: 25px; top: 20px; font-size: 10px; border-collapse: collapse; }}
  .meta td {{ padding: 2px 6px; border: 1px solid #ddd; }} .meta .lbl {{ font-weight: bold; color: {DK2}; }}
  .qor-tbl {{ width: 100%; border-collapse: collapse; font-size: 9px; margin-top: 8px; }}
  .qor-tbl th, .qor-tbl td {{ border: 1px solid #ccc; padding: 3px 5px; text-align: center; }}
  .qor-tbl .cat-row td {{ background: {DK2}; color: white; font-weight: bold; font-size: 8px; }}
  .qor-tbl .hdr-row th {{ background: {LT2}; color: {DK2}; font-size: 8px; }}
  .qor-tbl .stage {{ font-weight: bold; text-align: left; }}
  .neg {{ color: red; font-weight: bold; }}
  .two-col {{ display: flex; gap: 20px; }}
  .left-col {{ flex: 1.2; }} .right-col {{ flex: 1; }}
  .qor-mini {{ border-collapse: collapse; font-size: 10px; width: 100%; }}
  .qor-mini td {{ padding: 3px 6px; border: 1px solid #eee; }}
  .qor-mini td:first-child {{ font-weight: bold; color: {DK2}; background: #f8f9fa; }}
  .clock-tbl {{ border-collapse: collapse; font-size: 9px; margin-top: 10px; width: 100%; }}
  .clock-tbl th {{ background: {ACCENT1}; color: white; padding: 4px 6px; }}
  .clock-tbl td {{ padding: 3px 6px; border: 1px solid #ddd; }}
  .summary-text {{ font-size: 10px; color: {DK2}; min-height: 60px; }}
  .img-grid {{ display: grid; grid-template-columns: 1fr 1fr; gap: 8px; margin-top: 10px; }}
  .img-ph {{ background: {LT2}; border: 1px solid {ACCENT3}; padding: 30px 10px; text-align: center;
            font-size: 9px; color: {ACCENT3}; border-radius: 4px; }}
  @media print {{ body {{ background: white; }} .slide {{ box-shadow: none; margin: 0; }} }}
</style></head>
<body>
  <!-- Slide 1: Summary -->
  <div class="slide">
    <div class="slide-hdr"><span class="stage-title">Run : {data['run_name']} Summary</span></div>
    {meta}
    <table class="qor-tbl">
      <tr class="cat-row"><td colspan="19" style="background:{ACCENT1};">Run Name : {data['run_name']}</td></tr>
      {qor_header}
      {qor_rows}
    </table>
    <div style="margin-top:12px;">
      <strong style="color:{DK2};">Summary / Issues / Jira Tickets:</strong>
      <div style="min-height:40px;border:1px solid #eee;padding:6px;margin-top:4px;font-size:10px;color:#666;">
      </div>
    </div>
  </div>
  {stage_slides}
  <div style="text-align:center;padding:15px;font-size:11px;color:#999;">
    Generated by CBflow AutoPPT | {data['timestamp']}
  </div>
</body></html>"""

    with open(output_path, 'w') as f:
        f.write(html)


# ══════════════════════════════════════════════════════════════════════════════
# CLI
# ══════════════════════════════════════════════════════════════════════════════

def cmd_autoppt(args: argparse.Namespace) -> int:
    """Generate PowerPoint/HTML run summary."""
    run_dir = os.getcwd()

    if not os.path.exists(os.path.join(run_dir, '.stamps')) and \
       not any(Path(run_dir).glob('.race_*.db')) and \
       not os.path.exists(os.path.join(run_dir, '.run.cbflow.tcl')):
        print("ERROR: Not a CBflow run directory.")
        return 1

    print("Collecting run data...")
    data = collect_ppt_data(run_dir)

    output_dir = os.path.join(run_dir, 'reports')
    os.makedirs(output_dir, exist_ok=True)
    timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')

    if args.format == 'pptx' and not HAS_PPTX:
        print("WARNING: python-pptx not installed. Install with: pip install python-pptx")
        print("         Falling back to HTML format.")
        args.format = 'html'

    if args.format == 'pptx':
        output_file = args.output or os.path.join(output_dir, f'{data["run_name"]}_summary_{timestamp}.pptx')
        generate_pptx(data, output_file)
    else:
        output_file = args.output or os.path.join(output_dir, f'{data["run_name"]}_summary_{timestamp}.html')
        generate_html_slides(data, output_file)

    print(f"\nPD Run Summary generated: {output_file}")
    print(f"  Design     : {data['design_name']}")
    print(f"  Flow       : {data['flow_type']}")
    print(f"  Stages     : {len([s for s in data['stage_qor'].values() if any(s.values())])}/{len(STAGE_MAP)}")
    print(f"  Clock nets : {len(data['clock_table'])}")
    return 0


def create_parser(subparsers=None):
    """Create autoppt command parser."""
    _fmt = argparse.RawDescriptionHelpFormatter
    desc = """Generate PD Run Summary PPT/HTML matching the standard template.

Slides:
  1. Summary — 19-column QoR table (Setup/Hold/IO WNS/TNS/FEP, Design Stats,
               Clock Stats, Routability, Power, Runtime) + image placeholders
  2. Place   — Timing QoR + Placement/Density/Congestion/Pin images
  3. CTS     — Timing QoR + Clock skew table (Clock Name/ID/Skew)
  4. ctsOpt  — Timing QoR + Hold cells/Density/Congestion/Pin images
  5. routeOpt— Timing QoR + DRC/Shorts/Density/Pin images

Data is auto-collected from work/<flow>/<node>/reports/ directories.

Examples:
  cbflow run autoppt                      Generate HTML summary
  cbflow run autoppt --format pptx        Generate PowerPoint (requires python-pptx)
  cbflow run autoppt -o my_summary.html   Custom output path"""

    if subparsers:
        parser = subparsers.add_parser('autoppt', help='Generate PD run summary PPT/HTML',
            formatter_class=_fmt, description=desc)
    else:
        parser = argparse.ArgumentParser(prog='cbflow run autoppt',
            formatter_class=_fmt, description=desc)

    parser.add_argument('--format', '-f', choices=['pptx', 'html'], default='html',
                        help='Output format (default: html, pptx requires python-pptx)')
    parser.add_argument('--output', '-o', help='Output file path')
    return parser
