#!/usr/bin/env python3
"""Generate CBflow Feature Presentation (PowerPoint)"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Color Palette ──────────────────────────────────────────────────────────
DARK_BG      = RGBColor(0x1A, 0x1A, 0x2E)   # Deep navy
ACCENT_BLUE  = RGBColor(0x00, 0x96, 0xD6)   # Bright blue
ACCENT_TEAL  = RGBColor(0x00, 0xB4, 0xAB)   # Teal
ACCENT_GREEN = RGBColor(0x2E, 0xCC, 0x71)   # Green
ACCENT_ORANGE= RGBColor(0xF3, 0x9C, 0x12)   # Orange
ACCENT_RED   = RGBColor(0xE7, 0x4C, 0x3C)   # Red
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY   = RGBColor(0xCC, 0xCC, 0xCC)
MED_GRAY     = RGBColor(0x99, 0x99, 0x99)
CARD_BG      = RGBColor(0x24, 0x24, 0x3E)   # Slightly lighter navy
SECTION_BG   = RGBColor(0x16, 0x21, 0x3E)   # Section divider bg

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

W = prs.slide_width
H = prs.slide_height

# ── Helper Functions ───────────────────────────────────────────────────────

def add_bg(slide, color=DARK_BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, left, top, width, height, color, alpha=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    if alpha is not None:
        from lxml import etree
        solidFill = shape.fill._fill
        srgb = solidFill.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}srgbClr')
        if srgb is not None:
            alpha_elem = etree.SubElement(srgb, '{http://schemas.openxmlformats.org/drawingml/2006/main}alpha')
            alpha_elem.set('val', str(int(alpha * 1000)))
    return shape

def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_bullet_frame(slide, left, top, width, height, items, font_size=16,
                     color=WHITE, bullet_color=ACCENT_BLUE, spacing=Pt(6)):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_after = spacing
        # Bullet char
        run_b = p.add_run()
        run_b.text = "\u25B8 "
        run_b.font.size = Pt(font_size)
        run_b.font.color.rgb = bullet_color
        run_b.font.name = "Calibri"
        # Text
        run_t = p.add_run()
        run_t.text = item
        run_t.font.size = Pt(font_size)
        run_t.font.color.rgb = color
        run_t.font.name = "Calibri"
    return txBox

def add_card(slide, left, top, width, height, title, body_items,
             title_color=ACCENT_BLUE, card_color=CARD_BG):
    add_rect(slide, left, top, width, height, card_color)
    # accent bar
    add_rect(slide, left, top, Inches(0.06), height, title_color)
    add_text_box(slide, left + Inches(0.25), top + Inches(0.1), width - Inches(0.35),
                 Inches(0.4), title, font_size=16, color=title_color, bold=True)
    add_bullet_frame(slide, left + Inches(0.25), top + Inches(0.5),
                     width - Inches(0.35), height - Inches(0.6),
                     body_items, font_size=13, color=LIGHT_GRAY, bullet_color=MED_GRAY)

def add_accent_bar(slide, left, top, width, color=ACCENT_BLUE):
    add_rect(slide, left, top, width, Inches(0.05), color)

def section_number_badge(slide, left, top, number, color=ACCENT_BLUE):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(0.55), Inches(0.55))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.text = str(number)
    p.font.size = Pt(22)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = "Calibri"
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 1: TITLE SLIDE
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
add_bg(slide)
add_rect(slide, Inches(0), Inches(0), W, Inches(0.08), ACCENT_BLUE)
add_rect(slide, Inches(0), Inches(7.42), W, Inches(0.08), ACCENT_TEAL)

add_text_box(slide, Inches(1), Inches(1.5), Inches(11), Inches(1),
             "CBflow", font_size=64, color=WHITE, bold=True)
add_accent_bar(slide, Inches(1), Inches(2.6), Inches(2), ACCENT_BLUE)
add_text_box(slide, Inches(1), Inches(2.9), Inches(11), Inches(0.8),
             "Physical Design Automation Framework", font_size=32, color=ACCENT_TEAL)
add_text_box(slide, Inches(1), Inches(3.8), Inches(11), Inches(0.6),
             "Version 2.0.0  |  Complete Feature Overview", font_size=20, color=MED_GRAY)

add_text_box(slide, Inches(1), Inches(5.5), Inches(11), Inches(0.5),
             "Python  /  Bash  /  TCL  /  Makefile   |   13 Design Flows   |   Multi-Vendor EDA Support",
             font_size=16, color=LIGHT_GRAY, alignment=PP_ALIGN.LEFT)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 2: TABLE OF CONTENTS / AGENDA
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_rect(slide, Inches(0), Inches(0), W, Inches(0.06), ACCENT_BLUE)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.6),
             "Agenda", font_size=36, color=WHITE, bold=True)
add_accent_bar(slide, Inches(0.8), Inches(1.0), Inches(1.5), ACCENT_BLUE)

agenda_items = [
    ("1", "Architecture & Technology Stack", ACCENT_BLUE),
    ("2", "Supported Design Flows (13 Flows)", ACCENT_TEAL),
    ("3", "Multi-Vendor EDA Tool Support", ACCENT_GREEN),
    ("4", "CLI Command System (cbflow)", ACCENT_ORANGE),
    ("5", "Workspace & Run Management", ACCENT_BLUE),
    ("6", "Configuration Hierarchy", ACCENT_TEAL),
    ("7", "LSF & ML-Based Resource Optimization", ACCENT_GREEN),
    ("8", "Release & Version Management", ACCENT_ORANGE),
    ("9", "Plugin System & Extensibility", ACCENT_BLUE),
    ("10", "Metrics, Dashboard & Observability", ACCENT_TEAL),
    ("11", "Validation Framework", ACCENT_GREEN),
    ("12", "Advanced Features & Workflow", ACCENT_ORANGE),
]

col1 = agenda_items[:6]
col2 = agenda_items[6:]

for i, (num, label, color) in enumerate(col1):
    y = Inches(1.4) + Inches(i * 0.85)
    section_number_badge(slide, Inches(1.0), y, num, color)
    add_text_box(slide, Inches(1.7), y + Inches(0.08), Inches(4.5), Inches(0.45),
                 label, font_size=18, color=LIGHT_GRAY)

for i, (num, label, color) in enumerate(col2):
    y = Inches(1.4) + Inches(i * 0.85)
    section_number_badge(slide, Inches(7.0), y, num, color)
    add_text_box(slide, Inches(7.7), y + Inches(0.08), Inches(4.5), Inches(0.45),
                 label, font_size=18, color=LIGHT_GRAY)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 3: ARCHITECTURE & TECH STACK
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_rect(slide, Inches(0), Inches(0), W, Inches(0.06), ACCENT_BLUE)

section_number_badge(slide, Inches(0.8), Inches(0.35), "1", ACCENT_BLUE)
add_text_box(slide, Inches(1.5), Inches(0.35), Inches(10), Inches(0.55),
             "Architecture & Technology Stack", font_size=32, color=WHITE, bold=True)
add_accent_bar(slide, Inches(0.8), Inches(1.0), Inches(3), ACCENT_BLUE)

# 4 cards
cards = [
    ("Python 3.6+", [
        "20+ command modules (~400K lines)",
        "CLI handlers (workspace, run, flow)",
        "Makefile generator from TCL configs",
        "TCL config parser",
        "ML analytics for LSF optimization",
        "Flask-based web dashboard",
    ], ACCENT_BLUE),
    ("Bash / Zsh", [
        "Main CLI dispatcher (cbflow)",
        "5 sub-command routers",
        "Shell completion scripts",
        "Environment setup & validation",
        "Cross-shell compatibility",
    ], ACCENT_TEAL),
    ("TCL", [
        "EDA tool command scripts",
        "Flow/stage/subnode definitions",
        "Vendor-specific tool integration",
        "MMMC corner/mode configs",
        "13 flow node configurations",
    ], ACCENT_GREEN),
    ("Makefile / Make 3.81+", [
        "Dynamic per-run Makefile generation",
        "Stamp-based incremental execution",
        "Git-based component versioning",
        "Dependency graph management",
        "Release automation targets",
    ], ACCENT_ORANGE),
]

for i, (title, items, color) in enumerate(cards):
    left = Inches(0.5) + Inches(i * 3.15)
    add_card(slide, left, Inches(1.3), Inches(2.95), Inches(4.5), title, items, title_color=color)

add_text_box(slide, Inches(0.8), Inches(6.2), Inches(11.5), Inches(0.8),
             "System Requirements:  Git 2.0+  |  Python 3.6+  |  Make 3.81+  |  Bash/Zsh 4.0+",
             font_size=14, color=MED_GRAY, alignment=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 4: SUPPORTED DESIGN FLOWS (Part 1 - 7 flows)
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_rect(slide, Inches(0), Inches(0), W, Inches(0.06), ACCENT_TEAL)

section_number_badge(slide, Inches(0.8), Inches(0.35), "2", ACCENT_TEAL)
add_text_box(slide, Inches(1.5), Inches(0.35), Inches(10), Inches(0.55),
             "Supported Design Flows (1/2)", font_size=32, color=WHITE, bold=True)
add_accent_bar(slide, Inches(0.8), Inches(1.0), Inches(3), ACCENT_TEAL)

flows_1 = [
    ("SYNTH", "Logic Synthesis & Optimization", "RTL-to-gate synthesis, technology mapping, timing optimization",
     "inputs -> synthesis -> export_data -> release_data", ACCENT_BLUE),
    ("FP", "Floorplanning & Power Planning", "Die size, macro placement, power grid, pin assignment",
     "inputs -> floorplan -> export_data -> release_data", ACCENT_TEAL),
    ("PNR", "Place & Route", "Standard cell placement, CTS, routing, post-route optimization",
     "inputs -> place -> cts -> cts_opt -> route -> pro -> signoff -> export -> release", ACCENT_GREEN),
    ("STA", "Static Timing Analysis", "Sign-off timing, setup/hold checks, path analysis",
     "inputs -> timing -> export_data -> release_data", ACCENT_ORANGE),
    ("LEC", "Logic Equivalence Checking", "Formal verification of RTL vs netlist equivalence",
     "inputs -> compare -> export_data -> release_data", ACCENT_BLUE),
    ("EMIR", "EM/IR Drop Analysis", "Electromigration, IR drop, power integrity analysis",
     "inputs -> analysis -> export_data -> release_data", ACCENT_TEAL),
    ("PV", "Physical Verification", "DRC, LVS, ERC, PERC - full physical sign-off",
     "DRC, LVS, ERC, PERC stages", ACCENT_GREEN),
]

for i, (abbr, name, desc, stages, color) in enumerate(flows_1):
    y = Inches(1.3) + Inches(i * 0.82)
    add_rect(slide, Inches(0.5), y, Inches(12.3), Inches(0.72), CARD_BG)
    add_rect(slide, Inches(0.5), y, Inches(0.06), Inches(0.72), color)
    # Abbrev badge
    badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.75), y + Inches(0.1), Inches(0.9), Inches(0.5))
    badge.fill.solid()
    badge.fill.fore_color.rgb = color
    badge.line.fill.background()
    tf = badge.text_frame
    p = tf.paragraphs[0]
    p.text = abbr
    p.font.size = Pt(14)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    add_text_box(slide, Inches(1.85), y + Inches(0.05), Inches(3), Inches(0.35),
                 name, font_size=15, color=WHITE, bold=True)
    add_text_box(slide, Inches(1.85), y + Inches(0.35), Inches(4.5), Inches(0.35),
                 desc, font_size=11, color=MED_GRAY)
    add_text_box(slide, Inches(6.5), y + Inches(0.15), Inches(6), Inches(0.45),
                 stages, font_size=11, color=LIGHT_GRAY)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 5: SUPPORTED DESIGN FLOWS (Part 2 - 6 flows)
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_rect(slide, Inches(0), Inches(0), W, Inches(0.06), ACCENT_TEAL)

section_number_badge(slide, Inches(0.8), Inches(0.35), "2", ACCENT_TEAL)
add_text_box(slide, Inches(1.5), Inches(0.35), Inches(10), Inches(0.55),
             "Supported Design Flows (2/2)", font_size=32, color=WHITE, bold=True)
add_accent_bar(slide, Inches(0.8), Inches(1.0), Inches(3), ACCENT_TEAL)

flows_2 = [
    ("ECO", "Engineering Change Orders", "Post-signoff functional/timing ECOs, incremental optimization",
     "inputs -> eco -> export_data -> release_data", ACCENT_ORANGE),
    ("CLP", "Conformal Low Power", "UPF/CPF power intent verification, isolation/retention checks",
     "inputs -> lp_verify -> export_data -> release_data", ACCENT_BLUE),
    ("POPT", "Power Optimization", "Clock gating, power-aware optimization, leakage reduction",
     "inputs -> power_opt -> export_data -> release_data", ACCENT_TEAL),
    ("FCFP", "Full-Chip Floorplanning", "Top-level integration, hierarchical macro placement",
     "inputs -> floorplan -> export_data -> release_data", ACCENT_GREEN),
    ("PHYV", "Physical Verification (DFM)", "Design-for-manufacturing rule checks, yield optimization",
     "inputs -> verification -> export_data -> release_data", ACCENT_ORANGE),
    ("FCT", "Final Circuit Timing", "Final sign-off timing with extracted parasitics",
     "inputs -> timing -> export_data -> release_data", ACCENT_BLUE),
]

for i, (abbr, name, desc, stages, color) in enumerate(flows_2):
    y = Inches(1.3) + Inches(i * 0.82)
    add_rect(slide, Inches(0.5), y, Inches(12.3), Inches(0.72), CARD_BG)
    add_rect(slide, Inches(0.5), y, Inches(0.06), Inches(0.72), color)
    badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.75), y + Inches(0.1), Inches(0.9), Inches(0.5))
    badge.fill.solid()
    badge.fill.fore_color.rgb = color
    badge.line.fill.background()
    tf = badge.text_frame
    p = tf.paragraphs[0]
    p.text = abbr
    p.font.size = Pt(14)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    add_text_box(slide, Inches(1.85), y + Inches(0.05), Inches(3), Inches(0.35),
                 name, font_size=15, color=WHITE, bold=True)
    add_text_box(slide, Inches(1.85), y + Inches(0.35), Inches(4.5), Inches(0.35),
                 desc, font_size=11, color=MED_GRAY)
    add_text_box(slide, Inches(6.5), y + Inches(0.15), Inches(6), Inches(0.45),
                 stages, font_size=11, color=LIGHT_GRAY)

# Merged flows callout
add_rect(slide, Inches(0.5), Inches(6.3), Inches(12.3), Inches(0.8), CARD_BG)
add_rect(slide, Inches(0.5), Inches(6.3), Inches(12.3), Inches(0.06), ACCENT_ORANGE)
add_text_box(slide, Inches(0.8), Inches(6.4), Inches(5), Inches(0.35),
             "Merged Flow Support", font_size=16, color=ACCENT_ORANGE, bold=True)
add_text_box(slide, Inches(0.8), Inches(6.7), Inches(11.5), Inches(0.35),
             "Combine flows for license efficiency:  SYNTH_FP  |  SYNTH_FP_PNR  |  Custom combinations via merge_entry / merge_handoff / merge_parallel configs",
             font_size=12, color=LIGHT_GRAY)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 6: MULTI-VENDOR EDA TOOL SUPPORT
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_rect(slide, Inches(0), Inches(0), W, Inches(0.06), ACCENT_GREEN)

section_number_badge(slide, Inches(0.8), Inches(0.35), "3", ACCENT_GREEN)
add_text_box(slide, Inches(1.5), Inches(0.35), Inches(10), Inches(0.55),
             "Multi-Vendor EDA Tool Support", font_size=32, color=WHITE, bold=True)
add_accent_bar(slide, Inches(0.8), Inches(1.0), Inches(3), ACCENT_GREEN)

add_card(slide, Inches(0.5), Inches(1.3), Inches(3.8), Inches(5.3),
         "Synopsys Tools", [
             "Fusion Compiler (FC) - SYNTH, FP, PNR, ECO, FCFP",
             "PrimeTime (PT) - STA, POPT, FCT",
             "Formality - LEC",
             "RedHawk - EMIR analysis",
             "IC Validator (ICV) - PV (DRC/LVS/ERC/PERC)",
             "VC_LP - CLP low-power verification",
             "Power Compiler - POPT",
             "ICC - PNR (legacy)",
         ], title_color=ACCENT_BLUE)

add_card(slide, Inches(4.6), Inches(1.3), Inches(3.8), Inches(5.3),
         "Cadence Tools", [
             "Genus - Synthesis (SYNTH)",
             "Innovus - Floorplanning (FP) & PNR",
             "Tempus - Final Circuit Timing (FCT)",
             "Extensible to additional tools",
         ], title_color=ACCENT_TEAL)

add_card(slide, Inches(8.7), Inches(1.3), Inches(4.1), Inches(5.3),
         "Tool Management Features", [
             "Vendor/tool/version namespacing",
             "Per-flow tool configuration arrays",
             "Version-controlled command files",
             "Subnode handlers per tool variant",
             "Automatic tool path resolution",
             "Test mode: bypass EDA calls for dry runs",
             "Hot-swappable tool backends",
         ], title_color=ACCENT_GREEN)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 7: CLI COMMAND SYSTEM
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_rect(slide, Inches(0), Inches(0), W, Inches(0.06), ACCENT_ORANGE)

section_number_badge(slide, Inches(0.8), Inches(0.35), "4", ACCENT_ORANGE)
add_text_box(slide, Inches(1.5), Inches(0.35), Inches(10), Inches(0.55),
             "CLI Command System", font_size=32, color=WHITE, bold=True)
add_accent_bar(slide, Inches(0.8), Inches(1.0), Inches(3), ACCENT_ORANGE)

# Main command
add_rect(slide, Inches(0.5), Inches(1.25), Inches(12.3), Inches(0.55), CARD_BG)
add_text_box(slide, Inches(0.8), Inches(1.3), Inches(3), Inches(0.45),
             "cbflow", font_size=20, color=ACCENT_ORANGE, bold=True)
add_text_box(slide, Inches(2.5), Inches(1.3), Inches(10), Inches(0.45),
             "Main CLI dispatcher  -  routes to workspace / run / flow sub-commands",
             font_size=15, color=LIGHT_GRAY)

subcmds = [
    ("cbflow workspace", [
        "init - Initialize workspace for project/flow/block",
        "create - Create new run with user configuration",
        "status - Show workspace and run status/summary",
        "list-runs - List all runs in workspace",
        "validate - Validate workspace configuration",
        "clean - Clean workspace artifacts",
    ], ACCENT_BLUE),
    ("cbflow run", [
        "all - Execute all stages with optional validation",
        "stage <name> - Run a single stage",
        "status - Show run progress and details",
        "logs - View/tail logs with level filtering",
        "validate - Run validation checks on outputs",
        "retrace - Re-trace a completed stage",
        "lsf-status - LSF job status monitoring",
        "clean - Clean run artifacts",
    ], ACCENT_TEAL),
    ("cbflow flow", [
        "types - List all 13 supported flow types",
        "info <type> - Show flow details and stages",
        "stages <type> - List stages for a flow",
        "nodes - Show node dependency graph",
        "check - Verify tool installation",
        "config - Show/edit configuration",
        "plugin - Manage flow plugins",
        "metrics / dashboard / version / release",
    ], ACCENT_GREEN),
]

for i, (cmd, items, color) in enumerate(subcmds):
    left = Inches(0.5) + Inches(i * 4.15)
    add_card(slide, left, Inches(2.0), Inches(3.95), Inches(4.8), cmd, items, title_color=color)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 8: WORKSPACE & RUN MANAGEMENT
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_rect(slide, Inches(0), Inches(0), W, Inches(0.06), ACCENT_BLUE)

section_number_badge(slide, Inches(0.8), Inches(0.35), "5", ACCENT_BLUE)
add_text_box(slide, Inches(1.5), Inches(0.35), Inches(10), Inches(0.55),
             "Workspace & Run Management", font_size=32, color=WHITE, bold=True)
add_accent_bar(slide, Inches(0.8), Inches(1.0), Inches(3), ACCENT_BLUE)

add_card(slide, Inches(0.5), Inches(1.3), Inches(5.9), Inches(2.6),
         "Workspace Lifecycle", [
             "Project-scoped workspace initialization (project + flow + block)",
             "Automatic directory structure creation per flow type",
             "Multi-run support within single workspace",
             "Phase-based organization (P0, P1, P2, P3)",
             "Naming convention: P0_run_<FLOW>_run<N>",
             "Workspace-level validation and health checks",
         ], title_color=ACCENT_BLUE)

add_card(slide, Inches(6.9), Inches(1.3), Inches(5.9), Inches(2.6),
         "Run Execution Modes", [
             "Node mode: each stage as separate execution node (default)",
             "Flat mode: merged execution for license efficiency",
             "Stage-level subnodes: setup -> run -> validate -> finish",
             "Stamp-based incremental re-execution",
             "Dependency-aware stage ordering",
             "Parallel stage execution where possible",
         ], title_color=ACCENT_TEAL)

add_card(slide, Inches(0.5), Inches(4.2), Inches(5.9), Inches(2.8),
         "Run Monitoring & Control", [
             "Real-time run status with progress tracking",
             "Log viewing with level filtering (ERROR, WARN, INFO)",
             "Tail mode for live log streaming",
             "LSF job status integration",
             "Stage retrace for debugging",
             "Per-stage timing and resource metrics",
         ], title_color=ACCENT_GREEN)

add_card(slide, Inches(6.9), Inches(4.2), Inches(5.9), Inches(2.8),
         "Project Phase Management", [
             "P0 - Initial implementation and prototyping",
             "P1 - Design refinement and optimization",
             "P2 - Final implementation and sign-off",
             "P3 - Production and manufacturing release",
             "Exit milestones: FP_EXIT, PLACE_EXIT, CTS_EXIT, PRO_EXIT, BTO, MTO",
             "Phase-aware configuration loading",
         ], title_color=ACCENT_ORANGE)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 9: CONFIGURATION HIERARCHY
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_rect(slide, Inches(0), Inches(0), W, Inches(0.06), ACCENT_TEAL)

section_number_badge(slide, Inches(0.8), Inches(0.35), "6", ACCENT_TEAL)
add_text_box(slide, Inches(1.5), Inches(0.35), Inches(10), Inches(0.55),
             "Configuration Hierarchy", font_size=32, color=WHITE, bold=True)
add_accent_bar(slide, Inches(0.8), Inches(1.0), Inches(3), ACCENT_TEAL)

# Left: config layers
config_layers = [
    ("flow_config.tcl", "Master flow config - 8 sections covering project, flows, tools, runtime, I/O, validation, release, advanced", ACCENT_BLUE),
    ("dir_config.tcl", "Directory structure definitions per flow type - standardized folder hierarchy", ACCENT_TEAL),
    ("node_configs/*_config.tcl", "Per-flow configs (13 files) - stages, subnodes, dependencies, tools, timeouts", ACCENT_GREEN),
    ("lsf_config.tcl", "LSF queue management - 5 queue tiers with ML-based resource optimization", ACCENT_ORANGE),
    ("mmmc_config.tcl", "Multi-mode multi-corner analysis - corner/mode definitions, analysis setup", ACCENT_BLUE),
    ("validation_config.tcl", "Validation rules - exit criteria, naming conventions, file checks", ACCENT_TEAL),
    ("logo_config.tcl", "Branding and company information configuration", ACCENT_GREEN),
]

for i, (name, desc, color) in enumerate(config_layers):
    y = Inches(1.3) + Inches(i * 0.8)
    add_rect(slide, Inches(0.5), y, Inches(7.5), Inches(0.7), CARD_BG)
    add_rect(slide, Inches(0.5), y, Inches(0.06), Inches(0.7), color)
    add_text_box(slide, Inches(0.75), y + Inches(0.05), Inches(3.5), Inches(0.3),
                 name, font_size=13, color=color, bold=True)
    add_text_box(slide, Inches(0.75), y + Inches(0.35), Inches(7), Inches(0.3),
                 desc, font_size=11, color=MED_GRAY)

# Right: config scopes
add_card(slide, Inches(8.3), Inches(1.3), Inches(4.5), Inches(2.5),
         "Configuration Scopes", [
             "Project-level: ravendrive, india, phoenix",
             "Technology: gf_22nm, tsmc_7nm",
             "Flow-level: per-flow overrides",
             "Setup: common + project-specific",
             "User: run-time user overrides",
         ], title_color=ACCENT_TEAL)

add_card(slide, Inches(8.3), Inches(4.1), Inches(4.5), Inches(2.5),
         "Version-Controlled Configs", [
             "Config versions: v1.0.0, v1.0.1, workspace",
             "Edit-restricted configs for governance",
             "TCL-to-Python parser for tooling",
             "Dynamic Makefile generation from TCL",
             "Hierarchical variable resolution",
         ], title_color=ACCENT_ORANGE)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 10: LSF & ML-BASED RESOURCE OPTIMIZATION
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_rect(slide, Inches(0), Inches(0), W, Inches(0.06), ACCENT_GREEN)

section_number_badge(slide, Inches(0.8), Inches(0.35), "7", ACCENT_GREEN)
add_text_box(slide, Inches(1.5), Inches(0.35), Inches(10), Inches(0.55),
             "LSF & ML-Based Resource Optimization", font_size=32, color=WHITE, bold=True)
add_accent_bar(slide, Inches(0.8), Inches(1.0), Inches(3), ACCENT_GREEN)

# Queue tier table
queue_tiers = [
    ("S", "8 GB", "4 CPU", "2 hrs", "inputs, validation, quick checks", "1.0x", ACCENT_BLUE),
    ("M", "16 GB", "8 CPU", "4 hrs", "synthesis, floorplan, basic PNR", "2.0x", ACCENT_TEAL),
    ("L", "32 GB", "16 CPU", "8 hrs", "placement, routing, timing opt", "4.0x", ACCENT_GREEN),
    ("XL", "64 GB", "32 CPU", "12 hrs", "full-chip integration, hierarchical", "8.0x", ACCENT_ORANGE),
    ("Ultra", "128 GB", "64+ CPU", "24 hrs", "massive designs, extreme workloads", "16.0x", ACCENT_RED),
]

# Header
add_rect(slide, Inches(0.5), Inches(1.25), Inches(7.5), Inches(0.45), ACCENT_GREEN)
headers = ["Tier", "Memory", "CPU", "Limit", "Typical Workloads", "Cost"]
hx = [0.5, 1.2, 2.2, 3.1, 4.1, 7.0]
hw = [0.7, 1.0, 0.9, 1.0, 2.9, 0.8]
for j, h in enumerate(headers):
    add_text_box(slide, Inches(hx[j]), Inches(1.27), Inches(hw[j]), Inches(0.4),
                 h, font_size=12, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

for i, (tier, mem, cpu, limit, desc, cost, color) in enumerate(queue_tiers):
    y = Inches(1.75) + Inches(i * 0.52)
    bg_c = CARD_BG if i % 2 == 0 else DARK_BG
    add_rect(slide, Inches(0.5), y, Inches(7.5), Inches(0.48), bg_c)
    vals = [tier, mem, cpu, limit, desc, cost]
    for j, v in enumerate(vals):
        add_text_box(slide, Inches(hx[j]), y + Inches(0.05), Inches(hw[j]), Inches(0.38),
                     v, font_size=11, color=LIGHT_GRAY if j > 0 else color,
                     bold=(j == 0), alignment=PP_ALIGN.CENTER)

# ML Analytics card
add_card(slide, Inches(8.5), Inches(1.25), Inches(4.3), Inches(3.0),
         "ML Analytics Engine", [
             "Scikit-learn based resource prediction",
             "Historical job data analysis",
             "Dynamic queue recommendation",
             "Cost-performance optimization",
             "Automatic queue right-sizing",
         ], title_color=ACCENT_GREEN)

add_card(slide, Inches(8.5), Inches(4.5), Inches(4.3), Inches(2.5),
         "LSF Management Features", [
             "Job submission & queue management",
             "Max concurrency per queue tier",
             "Priority-based scheduling (20-50)",
             "Node-specific resource allocation",
             "Cost factor tracking & optimization",
         ], title_color=ACCENT_ORANGE)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 11: RELEASE & VERSION MANAGEMENT
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_rect(slide, Inches(0), Inches(0), W, Inches(0.06), ACCENT_ORANGE)

section_number_badge(slide, Inches(0.8), Inches(0.35), "8", ACCENT_ORANGE)
add_text_box(slide, Inches(1.5), Inches(0.35), Inches(10), Inches(0.55),
             "Release & Version Management", font_size=32, color=WHITE, bold=True)
add_accent_bar(slide, Inches(0.8), Inches(1.0), Inches(3), ACCENT_ORANGE)

add_card(slide, Inches(0.5), Inches(1.3), Inches(3.9), Inches(3.0),
         "Git-Based Versioning", [
             "Git branches replace directory copies",
             "60-90% disk space savings",
             "Full Git history & diff capabilities",
             "Workspace isolation via Git worktrees",
             "Component-level version tracking",
             "Semantic versioning (major.minor.patch)",
         ], title_color=ACCENT_BLUE)

add_card(slide, Inches(4.7), Inches(1.3), Inches(3.9), Inches(3.0),
         "Release System", [
             "Auto mode: quick releases from current state",
             "Config mode: controlled releases with JSON config",
             "MANIFEST.json with 41+ component versions",
             "Auto-generated CHANGELOG.md per release",
             "Symlinked release management (current -> vX.Y.Z)",
             "Release diff comparison",
         ], title_color=ACCENT_TEAL)

add_card(slide, Inches(8.9), Inches(1.3), Inches(3.9), Inches(3.0),
         "Component Management", [
             "41 tracked components per release",
             "Config flows, tech nodes, project configs",
             "Utility modules versioned independently",
             "Command modules per flow/vendor/tool",
             "Component checkout & promotion",
             "release_versions.tcl for TCL access",
         ], title_color=ACCENT_GREEN)

# Workflow
add_rect(slide, Inches(0.5), Inches(4.6), Inches(12.3), Inches(2.5), CARD_BG)
add_text_box(slide, Inches(0.8), Inches(4.7), Inches(11), Inches(0.4),
             "Version & Release Workflow", font_size=18, color=ACCENT_ORANGE, bold=True)

steps = [
    ("1. Create\nWorkspace", "git_create_workspace\nDIR=<component>"),
    ("2. Edit\nFiles", "Modify files in\nworktree workspace"),
    ("3. Commit\nVersion", "git_commit_version\nTYPE=patch|minor|major"),
    ("4. Promote\nVersion", "git_promote_version\nVERSION=vX.Y.Z"),
    ("5. Create\nRelease", "git_create_release\nTYPE=patch|minor|major"),
]

for i, (label, cmd) in enumerate(steps):
    x = Inches(0.8) + Inches(i * 2.4)
    badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(5.2), Inches(2.0), Inches(1.6))
    badge.fill.solid()
    badge.fill.fore_color.rgb = DARK_BG
    badge.line.color.rgb = ACCENT_ORANGE
    badge.line.width = Pt(1)
    tf = badge.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = label
    p.font.size = Pt(12)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.text = cmd
    p2.font.size = Pt(9)
    p2.font.color.rgb = MED_GRAY
    p2.alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    # Arrow
    if i < len(steps) - 1:
        add_text_box(slide, x + Inches(2.0), Inches(5.8), Inches(0.4), Inches(0.4),
                     "\u2192", font_size=24, color=ACCENT_ORANGE, alignment=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 12: PLUGIN SYSTEM & EXTENSIBILITY
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_rect(slide, Inches(0), Inches(0), W, Inches(0.06), ACCENT_BLUE)

section_number_badge(slide, Inches(0.8), Inches(0.35), "9", ACCENT_BLUE)
add_text_box(slide, Inches(1.5), Inches(0.35), Inches(10), Inches(0.55),
             "Plugin System & Extensibility", font_size=32, color=WHITE, bold=True)
add_accent_bar(slide, Inches(0.8), Inches(1.0), Inches(3), ACCENT_BLUE)

add_card(slide, Inches(0.5), Inches(1.3), Inches(5.9), Inches(3.0),
         "Plugin Architecture", [
             "Scaffold templates for new flow types",
             "Manifest-based plugin registration",
             "Hot-pluggable flow extensions",
             "Automatic integration with CLI system",
             "Plugin discovery and validation",
             "Version-aware plugin management",
         ], title_color=ACCENT_BLUE)

add_card(slide, Inches(6.8), Inches(1.3), Inches(5.9), Inches(3.0),
         "Extensibility Points", [
             "Custom flow type creation (new node configs)",
             "New vendor/tool backend integration",
             "Custom stage definitions and dependencies",
             "User-defined subnode handlers",
             "Custom validation rules and exit criteria",
             "Technology node library addition (gf_22nm, tsmc_7nm, ...)",
         ], title_color=ACCENT_TEAL)

add_card(slide, Inches(0.5), Inches(4.6), Inches(5.9), Inches(2.5),
         "Documentation System", [
             "6-section organized docs (Quick Start through Examples)",
             "Architecture & system design documentation",
             "Developer guide for contributors",
             "Configuration reference manual",
             "Makefile & Python API reference",
         ], title_color=ACCENT_GREEN)

add_card(slide, Inches(6.8), Inches(4.6), Inches(5.9), Inches(2.5),
         "Developer Tools", [
             "Shell completion scripts (Bash/Zsh)",
             "Test mode for dry-run validation",
             "Graph renderer for dependency visualization",
             "TCL config parser for Python tooling",
             "Log viewer with filtering and analysis",
         ], title_color=ACCENT_ORANGE)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 13: METRICS, DASHBOARD & OBSERVABILITY
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_rect(slide, Inches(0), Inches(0), W, Inches(0.06), ACCENT_TEAL)

section_number_badge(slide, Inches(0.8), Inches(0.35), "10", ACCENT_TEAL)
add_text_box(slide, Inches(1.5), Inches(0.35), Inches(10), Inches(0.55),
             "Metrics, Dashboard & Observability", font_size=32, color=WHITE, bold=True)
add_accent_bar(slide, Inches(0.8), Inches(1.0), Inches(3), ACCENT_TEAL)

add_card(slide, Inches(0.5), Inches(1.3), Inches(3.9), Inches(5.3),
         "Metrics Collection", [
             "Automated per-stage metrics gathering",
             "Timing metrics (WNS, TNS, FEP)",
             "Area and utilization tracking",
             "Power consumption analysis",
             "DRC/LVS violation counts",
             "Runtime and memory usage",
             "Metrics export (JSON/CSV)",
             "Historical trend analysis",
         ], title_color=ACCENT_BLUE)

add_card(slide, Inches(4.7), Inches(1.3), Inches(3.9), Inches(5.3),
         "Web Dashboard", [
             "Flask-based real-time UI",
             "Project-level visibility",
             "Run status overview",
             "Stage progress tracking",
             "Interactive metric charts",
             "Multi-project comparison",
             "Launched via: cbflow flow dashboard",
         ], title_color=ACCENT_TEAL)

add_card(slide, Inches(8.9), Inches(1.3), Inches(3.9), Inches(5.3),
         "Logging & Observability", [
             "Unified logging configuration",
             "Level-based filtering (ERROR/WARN/INFO/DEBUG)",
             "Per-stage log isolation",
             "Log viewer with tail mode",
             "Centralized log aggregation",
             "Structured log format",
             "Integrated with LSF job monitoring",
         ], title_color=ACCENT_GREEN)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 14: VALIDATION FRAMEWORK
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_rect(slide, Inches(0), Inches(0), W, Inches(0.06), ACCENT_GREEN)

section_number_badge(slide, Inches(0.8), Inches(0.35), "11", ACCENT_GREEN)
add_text_box(slide, Inches(1.5), Inches(0.35), Inches(10), Inches(0.55),
             "Validation Framework", font_size=32, color=WHITE, bold=True)
add_accent_bar(slide, Inches(0.8), Inches(1.0), Inches(3), ACCENT_GREEN)

add_card(slide, Inches(0.5), Inches(1.3), Inches(3.9), Inches(4.0),
         "Pre-Run Validation", [
             "Workspace configuration checks",
             "Tool installation verification",
             "Input file existence validation",
             "Naming convention enforcement",
             "Technology library availability",
             "License pre-check",
         ], title_color=ACCENT_BLUE)

add_card(slide, Inches(4.7), Inches(1.3), Inches(3.9), Inches(4.0),
         "In-Run Validation", [
             "Per-subnode validation (setup/run/validate/finish)",
             "Stage-level exit criteria checking",
             "Dependency satisfaction verification",
             "Output file generation checks",
             "Error detection and reporting",
             "Incremental validation via stamps",
         ], title_color=ACCENT_TEAL)

add_card(slide, Inches(8.9), Inches(1.3), Inches(3.9), Inches(4.0),
         "Post-Run Validation", [
             "Exit milestone verification",
             "Cross-stage consistency checks",
             "Metrics threshold validation",
             "Release readiness assessment",
             "Automated QoR reporting",
             "Sign-off checklist verification",
         ], title_color=ACCENT_GREEN)

# Milestones bar
add_rect(slide, Inches(0.5), Inches(5.6), Inches(12.3), Inches(1.4), CARD_BG)
add_text_box(slide, Inches(0.8), Inches(5.7), Inches(11), Inches(0.35),
             "Exit Milestones", font_size=16, color=ACCENT_GREEN, bold=True)

milestones = ["FP_EXIT", "PLACE_EXIT", "CTS_EXIT", "PRO_EXIT", "BTO", "MTO"]
for i, m in enumerate(milestones):
    x = Inches(0.8) + Inches(i * 2.0)
    badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(6.15), Inches(1.7), Inches(0.55))
    badge.fill.solid()
    badge.fill.fore_color.rgb = ACCENT_GREEN
    badge.line.fill.background()
    tf = badge.text_frame
    p = tf.paragraphs[0]
    p.text = m
    p.font.size = Pt(13)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    if i < len(milestones) - 1:
        add_text_box(slide, x + Inches(1.7), Inches(6.2), Inches(0.3), Inches(0.4),
                     "\u2192", font_size=18, color=ACCENT_GREEN)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 15: ADVANCED FEATURES
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_rect(slide, Inches(0), Inches(0), W, Inches(0.06), ACCENT_ORANGE)

section_number_badge(slide, Inches(0.8), Inches(0.35), "12", ACCENT_ORANGE)
add_text_box(slide, Inches(1.5), Inches(0.35), Inches(10), Inches(0.55),
             "Advanced Features & Workflow", font_size=32, color=WHITE, bold=True)
add_accent_bar(slide, Inches(0.8), Inches(1.0), Inches(3), ACCENT_ORANGE)

features = [
    ("Merged Flow Execution", [
        "Combine SYNTH + FP + PNR into single run",
        "merge_entry / merge_handoff / merge_parallel configs",
        "License-efficient execution mode",
    ], ACCENT_BLUE),
    ("Multi-Mode Multi-Corner (MMMC)", [
        "Corner and mode definition framework",
        "Cross-corner timing analysis",
        "MMMC-aware flow configuration",
    ], ACCENT_TEAL),
    ("Incremental Execution", [
        "Stamp-based stage tracking in .make/",
        "Skip completed stages on re-run",
        "Selective stage re-execution (retrace)",
    ], ACCENT_GREEN),
    ("Multi-Project Support", [
        "Master project registry (master_projects.json)",
        "Per-project tech/config scoping",
        "Projects: ravendrive, india, phoenix",
    ], ACCENT_ORANGE),
    ("Dependency Graph Engine", [
        "Visual graph rendering of flow DAGs",
        "Stage and subnode dependency tracking",
        "Automatic execution ordering",
    ], ACCENT_BLUE),
    ("Dynamic Makefile Generation", [
        "TCL config parsed to Python dicts",
        "Per-run Makefile auto-generated",
        "Target-per-stage with dependencies",
    ], ACCENT_TEAL),
]

for i, (title, items, color) in enumerate(features):
    col = i % 3
    row = i // 3
    left = Inches(0.5) + Inches(col * 4.15)
    top = Inches(1.3) + Inches(row * 2.9)
    add_card(slide, left, top, Inches(3.95), Inches(2.6), title, items, title_color=color)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 16: COMPLETE WORKFLOW EXAMPLE
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_rect(slide, Inches(0), Inches(0), W, Inches(0.06), ACCENT_BLUE)

add_text_box(slide, Inches(0.8), Inches(0.35), Inches(11), Inches(0.55),
             "End-to-End Workflow Example", font_size=32, color=WHITE, bold=True)
add_accent_bar(slide, Inches(0.8), Inches(1.0), Inches(3), ACCENT_BLUE)

workflow_steps = [
    ("Initialize", "cbflow workspace init\n--project ravendrive\n--flow PNR\n--block my_design", ACCENT_BLUE),
    ("Create Run", "cbflow workspace create\n--config user.tcl", ACCENT_TEAL),
    ("Execute", "cd P0_run_PNR_run1\ncbflow run all\n--validate", ACCENT_GREEN),
    ("Monitor", "cbflow run status\n--details\ncbflow run logs --tail 50", ACCENT_ORANGE),
    ("Validate", "cbflow run validate\ncbflow run lsf-status", ACCENT_BLUE),
    ("Release", "cbflow flow release\ncreate --type patch\n--desc 'Production'", ACCENT_TEAL),
]

for i, (label, cmd, color) in enumerate(workflow_steps):
    x = Inches(0.3) + Inches(i * 2.15)
    # Step number circle
    circ = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.7), Inches(1.4), Inches(0.5), Inches(0.5))
    circ.fill.solid()
    circ.fill.fore_color.rgb = color
    circ.line.fill.background()
    tf = circ.text_frame
    p = tf.paragraphs[0]
    p.text = str(i + 1)
    p.font.size = Pt(18)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    # Label
    add_text_box(slide, x + Inches(0.15), Inches(2.1), Inches(1.8), Inches(0.4),
                 label, font_size=16, color=color, bold=True, alignment=PP_ALIGN.CENTER)

    # Command box
    cmd_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x + Inches(0.05), Inches(2.6), Inches(1.95), Inches(2.0))
    cmd_box.fill.solid()
    cmd_box.fill.fore_color.rgb = CARD_BG
    cmd_box.line.color.rgb = color
    cmd_box.line.width = Pt(1)
    tf = cmd_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = cmd
    p.font.size = Pt(10)
    p.font.color.rgb = LIGHT_GRAY
    p.font.name = "Courier New"
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    # Arrow
    if i < len(workflow_steps) - 1:
        add_text_box(slide, x + Inches(2.0), Inches(1.5), Inches(0.2), Inches(0.35),
                     "\u2192", font_size=18, color=MED_GRAY)

# Bottom summary
add_rect(slide, Inches(0.5), Inches(5.2), Inches(12.3), Inches(1.8), CARD_BG)
add_text_box(slide, Inches(0.8), Inches(5.3), Inches(11.5), Inches(0.35),
             "Workspace Directory Structure", font_size=16, color=ACCENT_BLUE, bold=True)
add_text_box(slide, Inches(0.8), Inches(5.7), Inches(11.5), Inches(1.2),
             "workarea/\n"
             "  P0_run_PNR_run1/           # Phase_run_Flow_runN naming convention\n"
             "    .make/                    # Stamp files for incremental execution\n"
             "    logs/                     # Per-stage log files\n"
             "    reports/                  # Generated metrics and reports\n"
             "    outputs/                  # Stage output data",
             font_size=11, color=LIGHT_GRAY)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 17: KEY METRICS / SUMMARY
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_rect(slide, Inches(0), Inches(0), W, Inches(0.06), ACCENT_BLUE)
add_rect(slide, Inches(0), Inches(7.42), W, Inches(0.08), ACCENT_TEAL)

add_text_box(slide, Inches(0.8), Inches(0.35), Inches(11), Inches(0.55),
             "CBflow at a Glance", font_size=36, color=WHITE, bold=True)
add_accent_bar(slide, Inches(0.8), Inches(1.0), Inches(2), ACCENT_BLUE)

metrics = [
    ("13", "Design Flows", "Complete PD lifecycle coverage", ACCENT_BLUE),
    ("41+", "Components", "Version-tracked per release", ACCENT_TEAL),
    ("20+", "Python Modules", "~400K lines of automation", ACCENT_GREEN),
    ("5", "LSF Queue Tiers", "ML-optimized resource allocation", ACCENT_ORANGE),
    ("4", "Project Phases", "P0 through P3 lifecycle", ACCENT_BLUE),
    ("6", "Exit Milestones", "FP_EXIT through MTO", ACCENT_TEAL),
    ("2", "EDA Vendors", "Synopsys + Cadence support", ACCENT_GREEN),
    ("8", "Config Sections", "Hierarchical TCL configuration", ACCENT_ORANGE),
]

for i, (number, label, desc, color) in enumerate(metrics):
    col = i % 4
    row = i // 4
    x = Inches(0.5) + Inches(col * 3.15)
    y = Inches(1.4) + Inches(row * 2.8)

    add_rect(slide, x, y, Inches(2.95), Inches(2.4), CARD_BG)
    add_rect(slide, x, y, Inches(2.95), Inches(0.06), color)
    add_text_box(slide, x, y + Inches(0.3), Inches(2.95), Inches(0.8),
                 number, font_size=48, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x, y + Inches(1.1), Inches(2.95), Inches(0.4),
                 label, font_size=18, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x, y + Inches(1.5), Inches(2.95), Inches(0.5),
                 desc, font_size=12, color=MED_GRAY, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(0.8), Inches(6.8), Inches(11.5), Inches(0.5),
             "CBflow v2.0.0  -  Production-Grade Physical Design Automation",
             font_size=16, color=MED_GRAY, alignment=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════════════
# SAVE
# ═══════════════════════════════════════════════════════════════════════════
output_path = "/Users/vmerugu/projects/CBflow_clone/CBflow_Features_Presentation.pptx"
prs.save(output_path)
print(f"Presentation saved to: {output_path}")
print(f"Total slides: {len(prs.slides)}")
